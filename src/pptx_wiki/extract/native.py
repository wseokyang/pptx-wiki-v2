from __future__ import annotations

import hashlib
import html
import math
import re
from collections.abc import Iterable, Iterator
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models import BBox, DeckRecord, Element, SlideRecord, TableCell, TableData

_SAFE_FILENAME_RE = re.compile(r"[^A-Za-z0-9._-]+")


@dataclass(frozen=True, slots=True)
class _Affine:
    """Small 2-D affine matrix used to flatten nested PowerPoint groups.

    The six values represent::

        x' = a*x + c*y + e
        y' = b*x + d*y + f

    ``left @ right`` means "apply right, then left".
    """

    a: float = 1.0
    b: float = 0.0
    c: float = 0.0
    d: float = 1.0
    e: float = 0.0
    f: float = 0.0

    def __matmul__(self, other: _Affine) -> _Affine:
        return _Affine(
            a=self.a * other.a + self.c * other.b,
            b=self.b * other.a + self.d * other.b,
            c=self.a * other.c + self.c * other.d,
            d=self.b * other.c + self.d * other.d,
            e=self.a * other.e + self.c * other.f + self.e,
            f=self.b * other.e + self.d * other.f + self.f,
        )

    def point(self, x: float, y: float) -> tuple[float, float]:
        return self.a * x + self.c * y + self.e, self.b * x + self.d * y + self.f


def _translate(x: float, y: float) -> _Affine:
    return _Affine(e=x, f=y)


def _scale(x: float, y: float) -> _Affine:
    return _Affine(a=x, d=y)


def _rotate_clockwise(degrees: float) -> _Affine:
    # Office angles are clockwise. With the screen/slide y-axis pointing down,
    # this is the conventional rotation matrix with a positive angle.
    radians = math.radians(degrees)
    cosine = math.cos(radians)
    sine = math.sin(radians)
    return _Affine(a=cosine, b=sine, c=-sine, d=cosine)


def _content_matrix(
    parent: _Affine,
    *,
    x: int,
    y: int,
    width: int,
    height: int,
    rotation: float,
    flip_horizontal: bool,
    flip_vertical: bool,
) -> _Affine:
    """Map coordinates local to a shape's unrotated box into slide space."""

    center_x = x + width / 2
    center_y = y + height / 2
    return (
        parent
        @ _translate(center_x, center_y)
        @ _rotate_clockwise(rotation)
        @ _scale(-1.0 if flip_horizontal else 1.0, -1.0 if flip_vertical else 1.0)
        @ _translate(-width / 2, -height / 2)
    )


def _bbox_from_local(matrix: _Affine, x: float, y: float, width: float, height: float) -> BBox:
    points = (
        matrix.point(x, y),
        matrix.point(x + width, y),
        matrix.point(x + width, y + height),
        matrix.point(x, y + height),
    )
    min_x = math.floor(min(point[0] for point in points))
    min_y = math.floor(min(point[1] for point in points))
    max_x = math.ceil(max(point[0] for point in points))
    max_y = math.ceil(max(point[1] for point in points))
    return BBox(min_x, min_y, max(0, max_x - min_x), max(0, max_y - min_y))


def _int_or_zero(value: Any) -> int:
    try:
        return int(value or 0)
    except (TypeError, ValueError):
        return 0


def _shape_geometry(shape: Any) -> tuple[int, int, int, int]:
    return (
        _int_or_zero(getattr(shape, "left", 0)),
        _int_or_zero(getattr(shape, "top", 0)),
        _int_or_zero(getattr(shape, "width", 0)),
        _int_or_zero(getattr(shape, "height", 0)),
    )


def _shape_xfrm(shape: Any) -> Any | None:
    element = getattr(shape, "_element", None)
    if element is None:
        return None
    candidates = element.xpath(
        "./p:spPr/a:xfrm | ./p:grpSpPr/a:xfrm | "
        "./p:xfrm | ./p:graphicFrame/p:xfrm | ./p:pic/p:spPr/a:xfrm"
    )
    if candidates:
        return candidates[0]
    # A picture's `p:spPr` is already a direct child, so the first expression
    # normally finds it. Keep this local-name fallback for producer variants.
    candidates = element.xpath(".//*[local-name()='xfrm'][1]")
    return candidates[0] if candidates else None


def _shape_transform(shape: Any, parent: _Affine) -> tuple[_Affine, BBox, dict[str, Any]]:
    x, y, width, height = _shape_geometry(shape)
    xfrm = _shape_xfrm(shape)
    rotation = float(getattr(shape, "rotation", 0.0) or 0.0)
    flip_horizontal = bool(getattr(xfrm, "flipH", False))
    flip_vertical = bool(getattr(xfrm, "flipV", False))
    matrix = _content_matrix(
        parent,
        x=x,
        y=y,
        width=width,
        height=height,
        rotation=rotation,
        flip_horizontal=flip_horizontal,
        flip_vertical=flip_vertical,
    )
    bbox = _bbox_from_local(matrix, 0, 0, width, height)
    return matrix, bbox, {
        "rotation_degrees": rotation,
        "flip_horizontal": flip_horizontal,
        "flip_vertical": flip_vertical,
        "local_bbox_emu": {"x": x, "y": y, "width": width, "height": height},
    }


def _group_child_matrix(shape: Any, group_matrix: _Affine) -> _Affine:
    """Map the group's child coordinate system to slide coordinates."""

    _, _, width, height = _shape_geometry(shape)
    xfrm = _shape_xfrm(shape)
    if xfrm is None:
        return group_matrix
    child_offset = getattr(xfrm, "chOff", None)
    child_extent = getattr(xfrm, "chExt", None)
    if child_offset is None or child_extent is None:
        return group_matrix
    child_width = _int_or_zero(getattr(child_extent, "cx", 0))
    child_height = _int_or_zero(getattr(child_extent, "cy", 0))
    if child_width <= 0 or child_height <= 0:
        return group_matrix
    return (
        group_matrix
        @ _scale(width / child_width, height / child_height)
        @ _translate(
            -_int_or_zero(getattr(child_offset, "x", 0)),
            -_int_or_zero(getattr(child_offset, "y", 0)),
        )
    )


def _enum_name(value: Any) -> str | None:
    if value is None:
        return None
    name = getattr(value, "name", None)
    return str(name) if name is not None else str(value)


def _optional_float(value: Any) -> float | None:
    if value is None:
        return None
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _run_metadata(run: Any) -> dict[str, Any]:
    font = run.font
    size = getattr(font, "size", None)
    data: dict[str, Any] = {
        "text": str(getattr(run, "text", "")).replace("\v", "\n"),
        "font_name": getattr(font, "name", None),
        "font_size_pt": _optional_float(getattr(size, "pt", None)),
        "bold": getattr(font, "bold", None),
        "italic": getattr(font, "italic", None),
        "underline": getattr(font, "underline", None),
    }
    try:
        address = run.hyperlink.address
    except (AttributeError, KeyError, ValueError):
        address = None
    if address:
        data["hyperlink"] = address
    return {key: value for key, value in data.items() if value is not None}


def _bullet_metadata(paragraph: Any) -> dict[str, Any] | None:
    paragraph_properties = getattr(paragraph, "_p", None)
    paragraph_properties = (
        getattr(paragraph_properties, "pPr", None) if paragraph_properties is not None else None
    )
    if paragraph_properties is None:
        return None
    bullet_chars = paragraph_properties.xpath("./a:buChar/@char")
    auto_number = paragraph_properties.xpath("./a:buAutoNum/@type")
    if bullet_chars:
        return {"kind": "character", "value": bullet_chars[0]}
    if auto_number:
        return {"kind": "number", "value": auto_number[0]}
    if paragraph_properties.xpath("./a:buNone"):
        return {"kind": "none"}
    return None


def _text_frame_data(text_frame: Any) -> tuple[str, list[dict[str, Any]]]:
    paragraphs: list[dict[str, Any]] = []
    texts: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = str(paragraph.text).replace("\v", "\n")
        texts.append(text)
        data: dict[str, Any] = {
            "text": text,
            "level": int(getattr(paragraph, "level", 0) or 0),
            "alignment": _enum_name(getattr(paragraph, "alignment", None)),
            "runs": [_run_metadata(run) for run in paragraph.runs],
        }
        bullet = _bullet_metadata(paragraph)
        if bullet is not None:
            data["bullet"] = bullet
        paragraphs.append({key: value for key, value in data.items() if value is not None})
    return "\n".join(texts), paragraphs


def _shape_nonvisual_metadata(shape: Any) -> dict[str, Any]:
    element = getattr(shape, "_element", None)
    if element is None:
        return {}
    properties = element.xpath(".//*[local-name()='cNvPr'][1]")
    if not properties:
        return {}
    prop = properties[0]
    return {
        key: value
        for key, value in {
            "description": prop.get("descr"),
            "title": prop.get("title"),
            "hidden": prop.get("hidden") in {"1", "true", "True"},
        }.items()
        if value not in {None, "", False}
    }


def _shape_kind(shape: Any) -> str:
    if bool(getattr(shape, "has_table", False)):
        return "table"
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}:
        return "image"
    if bool(getattr(shape, "has_chart", False)) or shape_type == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if shape_type in {MSO_SHAPE_TYPE.DIAGRAM, MSO_SHAPE_TYPE.IGX_GRAPHIC}:
        return "diagram"
    if shape_type in {MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO}:
        return "media"
    if shape_type in {
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
    }:
        return "ole"
    if bool(getattr(shape, "has_text_frame", False)):
        return "text"
    return "shape"


def _safe_filename(name: str) -> str:
    safe = _SAFE_FILENAME_RE.sub("_", Path(name).name).strip("._")
    return safe or "image"


def _extract_image(shape: Any, assets_dir: Path | None) -> tuple[str | None, dict[str, Any]]:
    try:
        image = shape.image
        blob = image.blob
    except (AttributeError, KeyError, ValueError):
        return None, {"embedded_image_available": False}

    original_name = str(getattr(image, "filename", "image") or "image")
    digest = hashlib.sha256(blob).hexdigest()
    try:
        extension = str(image.ext or "bin").lstrip(".")
    except (AttributeError, KeyError, OSError, ValueError):
        # python-pptx asks Pillow for the format. That intentionally rejects
        # several Office-native formats (for example SVG/EMF), even though the
        # original relationship blob is still perfectly valid and valuable.
        extension = Path(original_name).suffix.lstrip(".") or "bin"
    metadata: dict[str, Any] = {
        "embedded_image_available": True,
        "image_sha256": digest,
        "image_filename": original_name,
        "image_extension": extension,
        "image_bytes": len(blob),
        "crop": {
            "left": float(getattr(shape, "crop_left", 0.0) or 0.0),
            "top": float(getattr(shape, "crop_top", 0.0) or 0.0),
            "right": float(getattr(shape, "crop_right", 0.0) or 0.0),
            "bottom": float(getattr(shape, "crop_bottom", 0.0) or 0.0),
        },
        # The raw blob does not include PowerPoint crop, rotation, overlays, or
        # grouped transforms. OCR should normally use the rendered slide ROI.
        "ocr_prefer_rendered_slide_crop": True,
    }
    if assets_dir is None:
        return None, metadata

    image_dir = assets_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)
    output_name = f"{digest[:16]}-{_safe_filename(original_name)}"
    if not output_name.lower().endswith(f".{extension.lower()}"):
        output_name = f"{output_name}.{extension}"
    output_path = image_dir / output_name
    if not output_path.exists():
        output_path.write_bytes(blob)
    return str(output_path), metadata


def _cell_bbox(
    matrix: _Affine,
    column_offsets: list[int],
    row_offsets: list[int],
    *,
    row: int,
    column: int,
    row_span: int,
    column_span: int,
    grid_scale_x: float,
    grid_scale_y: float,
) -> BBox:
    x1 = column_offsets[column] * grid_scale_x
    x2 = column_offsets[column + column_span] * grid_scale_x
    y1 = row_offsets[row] * grid_scale_y
    y2 = row_offsets[row + row_span] * grid_scale_y
    return _bbox_from_local(matrix, x1, y1, x2 - x1, y2 - y1)


def _bbox_dict(box: BBox) -> dict[str, int]:
    return {"x": box.x, "y": box.y, "width": box.width, "height": box.height}


def _table_html(data: TableData) -> str:
    by_row: dict[int, list[TableCell]] = {}
    for cell in data.cells:
        by_row.setdefault(cell.row, []).append(cell)
    lines = ["<table>", "  <tbody>"]
    for row in range(data.rows):
        lines.append("    <tr>")
        for cell in sorted(by_row.get(row, []), key=lambda item: item.column):
            attrs: list[str] = []
            if cell.row_span > 1:
                attrs.append(f'rowspan="{cell.row_span}"')
            if cell.column_span > 1:
                attrs.append(f'colspan="{cell.column_span}"')
            suffix = f" {' '.join(attrs)}" if attrs else ""
            value = html.escape(cell.text).replace("\n", "<br>")
            lines.append(f"      <td{suffix}>{value}</td>")
        lines.append("    </tr>")
    lines.extend(["  </tbody>", "</table>"])
    return "\n".join(lines)


def _extract_table(shape: Any, matrix: _Affine, slide_width: int, slide_height: int) -> tuple[TableData, dict[str, Any], str]:
    table = shape.table
    row_count = len(table.rows)
    column_count = len(table.columns)
    column_widths = [_int_or_zero(column.width) for column in table.columns]
    row_heights = [_int_or_zero(row.height) for row in table.rows]
    column_offsets = [0]
    row_offsets = [0]
    for width in column_widths:
        column_offsets.append(column_offsets[-1] + width)
    for height in row_heights:
        row_offsets.append(row_offsets[-1] + height)

    _, _, shape_width, shape_height = _shape_geometry(shape)
    grid_scale_x = shape_width / column_offsets[-1] if column_offsets[-1] > 0 else 1.0
    grid_scale_y = shape_height / row_offsets[-1] if row_offsets[-1] > 0 else 1.0
    cells: list[TableCell] = []
    merge_count = 0
    for row_index in range(row_count):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            if bool(cell.is_spanned):
                continue
            is_explicit_merge_origin = bool(cell.is_merge_origin)
            row_span = int(cell.span_height) if is_explicit_merge_origin else 1
            column_span = int(cell.span_width) if is_explicit_merge_origin else 1
            # Guard malformed files whose span extends beyond the grid.
            row_span = max(1, min(row_span, row_count - row_index))
            column_span = max(1, min(column_span, column_count - column_index))
            if is_explicit_merge_origin:
                merge_count += 1
            text, paragraphs = _text_frame_data(cell.text_frame)
            bbox = _cell_bbox(
                matrix,
                column_offsets,
                row_offsets,
                row=row_index,
                column=column_index,
                row_span=row_span,
                column_span=column_span,
                grid_scale_x=grid_scale_x,
                grid_scale_y=grid_scale_y,
            )
            cells.append(
                TableCell(
                    row=row_index,
                    column=column_index,
                    text=text,
                    row_span=row_span,
                    column_span=column_span,
                    # TableData contains only visible/canonical cells; shadow
                    # cells were skipped above. In the shared IR this flag
                    # distinguishes canonical cells from merge shadows, so a
                    # normal unmerged cell is also an origin.
                    is_merge_origin=True,
                    metadata={
                        "bbox_emu": _bbox_dict(bbox),
                        "bbox_normalized": bbox.normalized(slide_width, slide_height),
                        "paragraphs": paragraphs,
                        "is_explicit_merge_origin": is_explicit_merge_origin,
                    },
                )
            )

    data = TableData(rows=row_count, columns=column_count, cells=cells)
    metadata = {
        "column_widths_emu": column_widths,
        "row_heights_emu": row_heights,
        "merge_origin_count": merge_count,
        "style_options": {
            "first_row": bool(table.first_row),
            "first_column": bool(table.first_col),
            "last_row": bool(table.last_row),
            "last_column": bool(table.last_col),
            "horizontal_banding": bool(table.horz_banding),
            "vertical_banding": bool(table.vert_banding),
        },
        "native_boundary_authoritative": True,
        "ocr_policy": "never",
        "grid_to_shape_scale": {"x": grid_scale_x, "y": grid_scale_y},
    }
    return data, metadata, _table_html(data)


def _notes_text(slide: Any) -> str | None:
    if not bool(getattr(slide, "has_notes_slide", False)):
        return None
    text_frame = slide.notes_slide.notes_text_frame
    if text_frame is None:
        return None
    text, _ = _text_frame_data(text_frame)
    return text.strip() or None


def _core_properties(presentation: Any) -> dict[str, Any]:
    properties = presentation.core_properties
    keys = (
        "title",
        "subject",
        "author",
        "keywords",
        "comments",
        "category",
        "last_modified_by",
        "revision",
        "created",
        "modified",
    )
    data: dict[str, Any] = {}
    for key in keys:
        value = getattr(properties, key, None)
        if value is not None:
            data[key] = value.isoformat() if hasattr(value, "isoformat") else value
    return data


def _iter_shapes(shapes: Iterable[Any]) -> Iterator[tuple[int, Any]]:
    return iter(enumerate(shapes))


def extract_pptx(
    path: str | Path,
    *,
    assets_dir: str | Path | None = None,
    include_empty_shapes: bool = False,
    strict: bool = False,
) -> DeckRecord:
    """Extract authoritative native structure from a PPTX file.

    Text boxes and table cells are read from OOXML, not OCR. Group transforms
    are flattened into slide-space EMU boxes, while ``z_path`` and ``parent_id``
    retain the original hierarchy. Each native table remains one independent
    element, which lets the ROI stage treat two close tables as hard blockers
    instead of merging them into one image crop.

    Args:
        path: Input ``.pptx`` file.
        assets_dir: Optional directory for lossless extraction of embedded image
            blobs. OCR should still use a rendered slide crop when PowerPoint
            crop/rotation/overlays matter.
        include_empty_shapes: Include decorative shapes without native text.
        strict: Raise on the first malformed/unsupported shape. By default the
            extractor records a warning and continues with the remaining slide.
    """

    source_path = Path(path)
    if not source_path.is_file():
        raise FileNotFoundError(source_path)
    asset_root = Path(assets_dir) if assets_dir is not None else None
    presentation = Presentation(str(source_path))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    warnings: list[dict[str, Any]] = []
    slides: list[SlideRecord] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        elements: list[Element] = []
        next_z_index = 0

        def visit(
            shape: Any,
            parent_matrix: _Affine,
            z_path: tuple[int, ...],
            parent_id: str | None,
            ancestor_hidden: bool = False,
            *,
            _slide_number: int = slide_number,
            _elements: list[Element] = elements,
        ) -> None:
            nonlocal next_z_index
            shape_id = _int_or_zero(getattr(shape, "shape_id", 0))
            path_token = ".".join(f"{index:03d}" for index in z_path)
            element_id = f"s{_slide_number:04d}-p{path_token}-id{shape_id}"
            try:
                matrix, bbox, transform_metadata = _shape_transform(shape, parent_matrix)
                kind = _shape_kind(shape)
                nonvisual_metadata = _shape_nonvisual_metadata(shape)
                own_hidden = bool(nonvisual_metadata.get("hidden", False))
                effective_hidden = ancestor_hidden or own_hidden
                metadata: dict[str, Any] = {
                    **transform_metadata,
                    **nonvisual_metadata,
                    "shape_id": shape_id,
                    "shape_type": _enum_name(getattr(shape, "shape_type", None)),
                    "z_path": list(z_path),
                    "bbox_normalized": bbox.normalized(slide_width, slide_height),
                    "source_layer": "slide",
                    "effective_hidden": effective_hidden,
                }
                element = Element(
                    id=element_id,
                    slide_number=_slide_number,
                    kind=kind,
                    bbox=bbox,
                    z_index=next_z_index,
                    source="native",
                    name=getattr(shape, "name", None),
                    parent_id=parent_id,
                    metadata=metadata,
                )
                next_z_index += 1

                if kind == "group":
                    element.metadata.update({"container": True, "ocr_policy": "children_only"})
                    _elements.append(element)
                    child_parent = _group_child_matrix(shape, matrix)
                    for child_index, child in _iter_shapes(shape.shapes):
                        visit(
                            child,
                            child_parent,
                            (*z_path, child_index),
                            element_id,
                            effective_hidden,
                        )
                    return

                if kind == "table":
                    table_data, table_metadata, table_html = _extract_table(
                        shape, matrix, slide_width, slide_height
                    )
                    element.table = table_data
                    element.html = table_html
                    element.markdown = table_html  # merged cells require HTML, not pipe-table syntax
                    element.text = "\n".join(cell.text for cell in table_data.cells if cell.text)
                    element.metadata.update(table_metadata)
                    _elements.append(element)
                    return

                if bool(getattr(shape, "has_text_frame", False)):
                    text, paragraphs = _text_frame_data(shape.text_frame)
                    element.text = text
                    element.metadata["paragraphs"] = paragraphs
                    element.metadata["ocr_policy"] = "never"
                    try:
                        placeholder = shape.placeholder_format
                    except (AttributeError, ValueError):
                        placeholder = None
                    if placeholder is not None:
                        element.metadata["placeholder_type"] = _enum_name(placeholder.type)

                if kind in {"image", "ole"}:
                    asset_path, image_metadata = _extract_image(shape, asset_root)
                    element.asset_path = asset_path
                    element.metadata.update(image_metadata)

                if kind in {"image", "chart", "diagram", "media", "ole"}:
                    element.metadata.update(
                        {
                            "ocr_policy": "rendered_roi",
                            "native_boundary_authoritative": True,
                            "needs_visual_parse": True,
                        }
                    )

                if (
                    include_empty_shapes
                    or kind in {"image", "chart", "diagram", "media", "ole"}
                    or bool((element.text or "").strip())
                ):
                    _elements.append(element)
            except Exception as exc:
                if strict:
                    raise
                warnings.append(
                    {
                        "slide_number": _slide_number,
                        "element_id": element_id,
                        "shape_name": getattr(shape, "name", None),
                        "error": f"{type(exc).__name__}: {exc}",
                    }
                )

        for top_index, top_shape in _iter_shapes(slide.shapes):
            visit(top_shape, _Affine(), (top_index,), None)

        title_shape = getattr(slide.shapes, "title", None)
        title = None
        if title_shape is not None and bool(getattr(title_shape, "has_text_frame", False)):
            title = str(title_shape.text).strip() or None
        show_value = slide._element.get("show")
        slides.append(
            SlideRecord(
                number=slide_number,
                width=slide_width,
                height=slide_height,
                title=title,
                notes=_notes_text(slide),
                elements=elements,
                metadata={
                    "hidden": show_value in {"0", "false", "False"},
                    "native_element_count": len(elements),
                    "source_part": str(slide.part.partname),
                },
            )
        )

    return DeckRecord(
        source_path=str(source_path),
        slide_width=slide_width,
        slide_height=slide_height,
        slides=slides,
        metadata={
            "extractor": "python-pptx-native",
            "core_properties": _core_properties(presentation),
            "warnings": warnings,
            "warning_count": len(warnings),
        },
    )
