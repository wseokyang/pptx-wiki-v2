from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_wiki.models import BBox
from pptx_wiki.roi import non_overlapping_padded_bbox


def _bbox(shape) -> BBox:
    return BBox(int(shape.left), int(shape.top), int(shape.width), int(shape.height))


def test_fixture_contains_two_distinct_tables_with_one_point_gap(complex_pptx) -> None:
    path, spec = complex_pptx
    deck = Presentation(path)
    shapes = {shape.name: shape for shape in deck.slides[0].shapes}

    table_a = shapes[spec.table_a_name]
    table_b = shapes[spec.table_b_name]

    assert table_a is not table_b
    assert table_a.has_table and table_b.has_table
    assert table_b.top - (table_a.top + table_a.height) == spec.table_gap_emu
    assert "표 A" in table_a.table.cell(0, 0).text
    assert table_a.table.cell(0, 0).is_merge_origin
    assert table_a.table.cell(0, 0).span_width == 4
    assert table_b.table.cell(0, 0).is_merge_origin
    assert table_b.table.cell(0, 0).span_height == 2


def test_close_table_rois_remain_disjoint(complex_pptx) -> None:
    path, spec = complex_pptx
    deck = Presentation(path)
    shapes = {shape.name: shape for shape in deck.slides[0].shapes}
    table_a = _bbox(shapes[spec.table_a_name])
    table_b = _bbox(shapes[spec.table_b_name])

    crop_a = non_overlapping_padded_bbox(
        table_a,
        [table_b],
        slide_width=deck.slide_width,
        slide_height=deck.slide_height,
        padding_ratio=0.05,
    )
    crop_b = non_overlapping_padded_bbox(
        table_b,
        [table_a],
        slide_width=deck.slide_width,
        slide_height=deck.slide_height,
        padding_ratio=0.05,
    )

    assert crop_a.y2 <= crop_b.y
    assert crop_a.y <= table_a.y and crop_a.y2 >= table_a.y2
    assert crop_b.y <= table_b.y and crop_b.y2 >= table_b.y2


def test_fixture_has_raster_table_and_group_children(complex_pptx) -> None:
    path, spec = complex_pptx
    slide = Presentation(path).slides[0]
    shapes = {shape.name: shape for shape in slide.shapes}

    picture = shapes[spec.image_table_name]
    assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
    assert not picture.has_table

    group = shapes[spec.group_name]
    assert group.shape_type == MSO_SHAPE_TYPE.GROUP
    children = {shape.name: shape for shape in group.shapes}
    assert set(children) == {"GROUP_BACKGROUND", "GROUP_HEADING", "GROUP_VALUE"}
    assert children["GROUP_HEADING"].text == "그룹 내부 핵심 지표"
    assert "9,876만원" in children["GROUP_VALUE"].text
