from __future__ import annotations

from hashlib import sha256
import json
from pathlib import Path
import re
import shutil
from typing import Any, Mapping, Sequence

from pptx import Presentation
from pptx.util import Inches
import pytest

from pptx_wiki.collection import CollectionConfig, run_collection
from pptx_wiki.integration import (
    IntegrationConfig,
    QUALIFIED_CITATION_RE,
    validate_integrated_artifact,
)
from pptx_wiki.semantic import SemanticConfig


_LOCAL_CITATION_RE = re.compile(r"\[slide-\d+#[^\]\s#]+\]")


def _write_multi_pr_reliability_deck(path: Path) -> Path:
    """Create a deck whose two PR numbers occur only in a body shape."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Semiconductor Package Reliability Analysis"
    body = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(8.5), Inches(4.5)
    )
    body.name = "ANALYSIS_BODY"
    body.text_frame.text = (
        "Request PR-10001\n"
        "Package: FC-BGA\n"
        "Lot No.: LOT-A7\n"
        "Test: Temperature Cycling Test\n"
        "Condition: -55 °C to 125 °C, 1000 cycles\n"
        "Result: 0/77 fail\n"
        "Request PR-10002\n"
        "Package: FC-BGA\n"
        "Test: HAST\n"
        "Condition: 130 °C / 85% RH / 96 h\n"
        "Result: 2/77 fail\n"
        "Failure mode: delamination"
    )
    presentation.save(path)
    return path


def _allowed_citations(prompt: str, pattern: re.Pattern[str]) -> tuple[str, ...]:
    line = next(
        value for value in prompt.splitlines() if value.startswith("Allowed citations:")
    )
    return tuple(dict.fromkeys(match.group(0) for match in pattern.finditer(line)))


class _SemiconductorKgBackend:
    """Script the same JSON/Markdown exchanges expected from the user's API."""

    model = "scripted-semiconductor-kg"

    def __init__(self) -> None:
        self.calls: list[str] = []

    def complete(
        self,
        messages: Sequence[Mapping[str, str]],
        *,
        max_tokens: int,
        temperature: float,
    ) -> str:
        del max_tokens, temperature
        prompt = messages[-1]["content"]

        if "Write one concise wiki page" in prompt:
            self.calls.append("source-page")
            citation = _allowed_citations(prompt, _LOCAL_CITATION_RE)[0]
            return (
                "# Package reliability analysis\n\n"
                "- PR-10001 used FC-BGA lot LOT-A7 for Temperature Cycling Test "
                "at -55 °C to 125 °C for 1000 cycles, with 0/77 fail; "
                "PR-10002 used FC-BGA for HAST at 130 °C / 85% RH / 96 h, "
                f"with 2/77 fail and delamination. {citation}"
            )

        if "COLLECTION_ENTITY_TOPIC_DISCOVERY" in prompt:
            self.calls.append("integration-discovery")
            citation = _allowed_citations(prompt, QUALIFIED_CITATION_RE)[0]
            return json.dumps(
                {
                    "entities": [
                        {
                            "name": "FC-BGA",
                            "type": "product",
                            "description": f"Package under reliability analysis {citation}",
                            "aliases": [],
                            "citations": [citation],
                        }
                    ],
                    "topics": [
                        {
                            "title": "Package reliability results",
                            "description": "Request, test condition, result, and failure mode",
                            "citations": [citation],
                        }
                    ],
                }
            )

        if "COLLECTION_GROUNDED_PAGE" in prompt:
            self.calls.append("integration-page")
            citation = _allowed_citations(prompt, QUALIFIED_CITATION_RE)[0]
            return (
                "- PR-10001 Temperature Cycling Test result was 0/77 fail; "
                "PR-10002 HAST result was 2/77 fail with delamination. "
                f"{citation}"
            )

        # The domain KG pass is intentionally a separate LLM exchange from the
        # generic topic planner.  Accept the stable marker while keeping this
        # fixture insensitive to explanatory prose around the JSON contract.
        if "COLLECTION_SEMICONDUCTOR_KG_EXTRACTION" in prompt:
            self.calls.append("domain-kg")
            citation = _allowed_citations(prompt, QUALIFIED_CITATION_RE)[0]
            return json.dumps(
                {
                    "entities": [
                        {
                            "type": "package",
                            "name": "FC-BGA",
                            "aliases": [],
                            "description": f"Package in the analysis {citation}",
                            "citations": [citation],
                        },
                        {
                            "type": "lot",
                            "name": "LOT-A7",
                            "aliases": [],
                            "description": f"Reliability sample lot {citation}",
                            "citations": [citation],
                        },
                        {
                            "type": "test_method",
                            "name": "Temperature Cycling Test",
                            "aliases": [],
                            "description": f"Reliability test method {citation}",
                            "citations": [citation],
                        },
                        {
                            "type": "test_method",
                            "name": "HAST",
                            "aliases": [],
                            "description": f"Reliability test method {citation}",
                            "citations": [citation],
                        },
                        {
                            "type": "failure_mode",
                            "name": "delamination",
                            "aliases": [],
                            "description": f"Observed failure mode {citation}",
                            "citations": [citation],
                        },
                    ],
                    "relationships": [
                        {
                            "predicate": "uses_package",
                            "subject": "PR-10001",
                            "object": "FC-BGA",
                            "assertion": "PR-10001 uses FC-BGA",
                            "description": f"Request uses package {citation}",
                            "citations": [citation],
                        },
                        {
                            "predicate": "has_lot",
                            "subject": "PR-10001",
                            "object": "LOT-A7",
                            "assertion": "PR-10001 has lot LOT-A7",
                            "description": f"Request has sample lot {citation}",
                            "citations": [citation],
                        },
                        {
                            "predicate": "underwent_test",
                            "subject": "PR-10001",
                            "object": "Temperature Cycling Test",
                            "assertion": (
                                "PR-10001 Temperature Cycling Test result was "
                                "0/77 fail"
                            ),
                            "description": f"Request specifies the test {citation}",
                            "citations": [citation],
                        },
                        {
                            "predicate": "underwent_test",
                            "subject": "PR-10002",
                            "object": "HAST",
                            "assertion": "PR-10002 HAST result was 2/77 fail",
                            "description": f"Request specifies the test {citation}",
                            "citations": [citation],
                        },
                        {
                            "predicate": "observed_failure",
                            "subject": "PR-10002",
                            "object": "delamination",
                            "assertion": "PR-10002 observed delamination",
                            "description": f"Request records a failure mode {citation}",
                            "citations": [citation],
                        },
                    ],
                }
            )

        raise AssertionError(f"unexpected backend prompt: {prompt[:160]}")


@pytest.fixture(scope="module")
def semiconductor_kg_collection(tmp_path_factory: pytest.TempPathFactory):
    root = tmp_path_factory.mktemp("semiconductor-kg")
    source = _write_multi_pr_reliability_deck(root / "analysis.pptx")
    backend = _SemiconductorKgBackend()
    result = run_collection(
        [source],
        root / "collection",
        semantic_backend=backend,
        integration_backend=backend,
        config=CollectionConfig(
            semantic=SemanticConfig(
                goal="Keep semiconductor package reliability requests and results.",
                coverage_policy="selected",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(
                goal="Build a grounded semiconductor package reliability KG.",
                repair_attempts=0,
            ),
            site_title="Semiconductor Reliability KG",
        ),
    )
    return result, backend


def _records(path: Path) -> list[dict[str, Any]]:
    return [
        json.loads(line)
        for line in path.read_text(encoding="utf-8").splitlines()
        if line.strip()
    ]


def _rewrite_integrated_jsonl(
    integrated: Path, name: str, records: Sequence[Mapping[str, Any]]
) -> None:
    text = "".join(
        json.dumps(value, ensure_ascii=False, sort_keys=True) + "\n"
        for value in records
    )
    (integrated / name).write_text(text, encoding="utf-8", newline="\n")
    manifest_path = integrated / "manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["files"][name] = {
        "sha256": sha256(text.encode("utf-8")).hexdigest(),
        "count": len(records),
    }
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, sort_keys=True, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )


def test_llm_builds_grounded_semiconductor_kg_and_quartz_links(
    semiconductor_kg_collection,
) -> None:
    result, backend = semiconductor_kg_collection

    # PR discovery is structure based: neither filename nor slide title owns it.
    assert result.pr_numbers == ("PR-10001", "PR-10002")
    assert backend.calls.count("domain-kg") == 1

    integrated = validate_integrated_artifact(result.integrated.output_dir)
    entities = integrated["entities.jsonl"]
    relationships = integrated["relationships.jsonl"]
    source_map_citations = {
        item["qualified_citation"] for item in integrated["source-map.jsonl"]
    }

    assert {entity["type"] for entity in entities} >= {
        "package",
        "lot",
        "test_method",
        "failure_mode",
    }

    entity_ids = {entity["id"] for entity in entities}
    for record in [*entities, *relationships]:
        assert record["citations"]
        assert set(record["citations"]) <= source_map_citations
        assert record["source_ids"] == [result.sources[0].source_id]
        assert set(record["pr_numbers"]) <= {"PR-10001", "PR-10002"}
    for relationship in relationships:
        for endpoint in (relationship["subject"], relationship["object"]):
            assert endpoint["kind"] in {"pr", "entity"}
            if endpoint["kind"] == "pr":
                assert endpoint["id"] in {"PR-10001", "PR-10002"}
            else:
                assert endpoint["id"] in entity_ids
        if relationship["subject"]["kind"] == "pr":
            assert relationship["pr_numbers"] == [relationship["subject"]["id"]]

    assertions = "\n".join(item["assertion"] for item in relationships)
    assert "0/77 fail" in assertions
    assert "2/77 fail" in assertions

    entity_pages = list((result.quartz.content_dir / "entities").glob("*.md"))
    pr_pages = list((result.quartz.content_dir / "prs").glob("*.md"))
    assert len(entity_pages) == len(entities)
    assert len(pr_pages) == 2
    rendered_entities = "\n".join(
        path.read_text(encoding="utf-8") for path in entity_pages
    )
    rendered_prs = "\n".join(path.read_text(encoding="utf-8") for path in pr_pages)
    assert "FC-BGA" in rendered_entities
    assert "[[prs/" in rendered_entities
    assert "[[entities/" in rendered_prs
    assert "0/77 fail" in rendered_prs
    assert "2/77 fail" in rendered_prs
    assert "[[evidence/" in rendered_entities + rendered_prs


def test_relationship_pr_endpoints_are_limited_to_detected_body_prs(
    semiconductor_kg_collection,
    tmp_path: Path,
) -> None:
    result, _ = semiconductor_kg_collection
    integrated = tmp_path / "integrated"
    shutil.copytree(result.integrated.output_dir, integrated)
    relationships = _records(integrated / "relationships.jsonl")
    relationships[0]["subject"] = {"kind": "pr", "id": "PR-99999"}
    relationships[0]["pr_numbers"] = ["PR-99999"]
    _rewrite_integrated_jsonl(integrated, "relationships.jsonl", relationships)

    with pytest.raises(ValueError, match="unknown.*PR|PR.*unknown|PR inventory"):
        validate_integrated_artifact(integrated)


def test_integrated_kg_rejects_hallucinated_citation(
    semiconductor_kg_collection,
    tmp_path: Path,
) -> None:
    result, _ = semiconductor_kg_collection
    integrated = tmp_path / "integrated"
    shutil.copytree(result.integrated.output_dir, integrated)
    relationships = _records(integrated / "relationships.jsonl")
    relationships[0]["citations"] = ["[@invented-source/slide-99#invented]"]
    _rewrite_integrated_jsonl(integrated, "relationships.jsonl", relationships)

    with pytest.raises(ValueError, match="unknown.*citation|citation.*unknown"):
        validate_integrated_artifact(integrated)


def test_integrated_kg_rejects_unknown_relationship_endpoint(
    semiconductor_kg_collection,
    tmp_path: Path,
) -> None:
    result, _ = semiconductor_kg_collection
    integrated = tmp_path / "integrated"
    shutil.copytree(result.integrated.output_dir, integrated)
    relationships = _records(integrated / "relationships.jsonl")
    relationships[0]["object"] = {"kind": "entity", "id": "unknown-entity"}
    _rewrite_integrated_jsonl(integrated, "relationships.jsonl", relationships)

    with pytest.raises(ValueError, match="unknown.*(?:endpoint|node)|(?:endpoint|node).*unknown"):
        validate_integrated_artifact(integrated)


def test_integrated_validator_keeps_legacy_artifacts_without_relationships(
    semiconductor_kg_collection,
    tmp_path: Path,
) -> None:
    result, _ = semiconductor_kg_collection
    integrated = tmp_path / "integrated"
    shutil.copytree(result.integrated.output_dir, integrated)
    (integrated / "relationships.jsonl").unlink()
    manifest_path = integrated / "manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    del manifest["files"]["relationships.jsonl"]
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, sort_keys=True, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )

    loaded = validate_integrated_artifact(integrated)
    assert loaded.get("relationships.jsonl", []) == []
