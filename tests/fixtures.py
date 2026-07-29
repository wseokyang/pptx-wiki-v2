"""Synthetic presentations used by extraction and export regression tests.

The fixture intentionally reproduces the failure mode that prompted this
project: two independent native tables are separated by only one point.  A
render-first OCR pipeline will commonly see one large grid, while PPTX shape
boundaries unambiguously identify two objects.
"""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class ComplexPptxSpec:
    """Facts tests may rely on without duplicating fixture coordinates."""

    table_gap_emu: int
    table_a_name: str = "NATIVE_TABLE_A"
    table_b_name: str = "NATIVE_TABLE_B"
    image_table_name: str = "IMAGE_TABLE_SCREENSHOT"
    group_name: str = "GROUPED_CALLOUT"


def _set_cell_text(cell, text: str, *, bold: bool = False) -> None:
    cell.text = text
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(10)
        run.font.bold = bold


def _image_table_stream() -> BytesIO:
    """Return a deterministic table-like PNG for picture/ROI tests."""

    image = Image.new("RGB", (800, 300), "white")
    draw = ImageDraw.Draw(image)
    navy = (32, 58, 91)
    pale_blue = (223, 236, 250)

    draw.rectangle((0, 0, 799, 299), outline=navy, width=5)
    draw.rectangle((0, 0, 799, 74), fill=pale_blue, outline=navy, width=4)
    for x in (260, 520):
        draw.line((x, 0, x, 300), fill=navy, width=4)
    for y in (75, 150, 225):
        draw.line((0, y, 800, y), fill=navy, width=4)

    labels = (
        (24, 26, "IMAGE-TABLE"),
        (284, 26, "CODE"),
        (544, 26, "VALUE"),
        (24, 101, "screenshot-row-1"),
        (284, 101, "IMG-A"),
        (544, 101, "1,234.50"),
        (24, 176, "screenshot-row-2"),
        (284, 176, "IMG-B"),
        (544, 176, "98.7%"),
        (24, 251, "flattened pixels only"),
    )
    for x, y, label in labels:
        draw.text((x, y), label, fill=(0, 0, 0))

    stream = BytesIO()
    image.save(stream, format="PNG", optimize=False)
    stream.seek(0)
    return stream


def build_complex_pptx(path: Path) -> ComplexPptxSpec:
    """Create a one-slide Korean PPTX containing extraction edge cases.

    Included cases:

    * several independent text boxes with a non-trivial z-order;
    * two native tables with a one-point vertical gap;
    * horizontal and vertical merged cells;
    * a table represented only as a raster picture;
    * a group with a background shape and two text-box children.
    """

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)

    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.3), Inches(0.48))
    title.name = "TITLE_TEXTBOX"
    title.text_frame.text = "분기 실적 — LLM Wiki 변환 회귀 테스트"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    lead = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(7.4), Inches(0.48))
    lead.name = "LEAD_TEXTBOX"
    lead.text_frame.text = "아래 두 표는 서로 다른 객체이며 간격은 1pt입니다."

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.25),
        Inches(0.78),
        Inches(1.25),
        Inches(0.42),
    )
    badge.name = "STATUS_BADGE"
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(226, 239, 218)
    badge.text_frame.text = "확정"

    # Close native table A: horizontal merge in the heading row.
    table_x = Inches(0.55)
    table_width = Inches(7.55)
    table_a_top = Inches(1.38)
    table_a_height = Inches(1.55)
    table_a_shape = slide.shapes.add_table(3, 4, table_x, table_a_top, table_width, table_a_height)
    table_a_shape.name = ComplexPptxSpec.table_a_name
    table_a = table_a_shape.table
    table_a.cell(0, 0).merge(table_a.cell(0, 3))
    _set_cell_text(table_a.cell(0, 0), "표 A — 국내 매출", bold=True)
    for column, value in enumerate(("제품", "수량", "단가", "매출")):
        _set_cell_text(table_a.cell(1, column), value, bold=True)
    for column, value in enumerate(("알파", "12", "1,250원", "15,000원")):
        _set_cell_text(table_a.cell(2, column), value)

    # Exactly one typographic point separates the shapes.  OCR dilation or
    # layout-model padding can easily erase this gap at ordinary render sizes.
    table_gap = Pt(1)
    table_b_top = table_a_top + table_a_height + table_gap
    table_b_height = Inches(1.55)
    table_b_shape = slide.shapes.add_table(3, 4, table_x, table_b_top, table_width, table_b_height)
    table_b_shape.name = ComplexPptxSpec.table_b_name
    table_b = table_b_shape.table
    table_b.cell(0, 0).merge(table_b.cell(1, 0))
    _set_cell_text(table_b.cell(0, 0), "표 B", bold=True)
    for column, value in enumerate(("지역", "목표", "실적"), start=1):
        _set_cell_text(table_b.cell(0, column), value, bold=True)
    for column, value in enumerate(("서울", "93.2%", "완료"), start=1):
        _set_cell_text(table_b.cell(1, column), value)
    for column, value in enumerate(("합계", "전국", "1,000", "987")):
        _set_cell_text(table_b.cell(2, column), value)

    footnote = slide.shapes.add_textbox(Inches(0.65), Inches(4.58), Inches(7.3), Inches(0.38))
    footnote.name = "TABLE_FOOTNOTE"
    footnote.text_frame.text = "주: 퍼센트와 쉼표가 포함된 숫자는 원문 그대로 보존해야 합니다."

    # Raster-only content ensures the extraction plan still creates an OCR/VLM
    # candidate without treating native neighbours as part of the same ROI.
    picture = slide.shapes.add_picture(
        _image_table_stream(),
        Inches(8.52),
        Inches(3.25),
        width=Inches(4.28),
        height=Inches(1.60),
    )
    picture.name = ComplexPptxSpec.image_table_name

    picture_label = slide.shapes.add_textbox(Inches(8.52), Inches(4.88), Inches(4.28), Inches(0.42))
    picture_label.name = "IMAGE_TABLE_CAPTION"
    picture_label.text_frame.text = "그림 1. 픽셀로만 존재하는 표"

    # Group children use slide coordinates; python-pptx recalculates group
    # extents as each child is appended.
    group = slide.shapes.add_group_shape()
    group.name = ComplexPptxSpec.group_name
    panel = group.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.52),
        Inches(1.38),
        Inches(4.28),
        Inches(1.48),
    )
    panel.name = "GROUP_BACKGROUND"
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(242, 242, 242)
    panel.text_frame.clear()
    group_heading = group.shapes.add_textbox(Inches(8.76), Inches(1.62), Inches(3.7), Inches(0.35))
    group_heading.name = "GROUP_HEADING"
    group_heading.text_frame.text = "그룹 내부 핵심 지표"
    group_heading.text_frame.paragraphs[0].runs[0].font.bold = True
    group_value = group.shapes.add_textbox(Inches(8.76), Inches(2.08), Inches(3.7), Inches(0.48))
    group_value.name = "GROUP_VALUE"
    group_value.text_frame.text = "영업이익 9,876만원 (+12.4%)"

    # A second slide checks stable slide numbering and repeated shape names.
    slide2 = deck.slides.add_slide(deck.slide_layouts[6])
    heading2 = slide2.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.55))
    heading2.name = "TITLE_TEXTBOX"
    heading2.text_frame.text = "부록 — 읽기 순서 안정성"
    left = slide2.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(5.8), Inches(1.0))
    left.name = "COLUMN_LEFT"
    left.text_frame.text = "왼쪽 열\n첫 번째 문단"
    right = slide2.shapes.add_textbox(Inches(6.85), Inches(1.25), Inches(5.8), Inches(1.0))
    right.name = "COLUMN_RIGHT"
    right.text_frame.text = "오른쪽 열\n두 번째 문단"

    deck.core_properties.title = "PPTX Wiki synthetic regression fixture"
    deck.core_properties.subject = "Adjacent tables, merged cells, pictures, and grouped shapes"
    deck.core_properties.author = "pptx-wiki tests"
    deck.save(path)
    return ComplexPptxSpec(table_gap_emu=int(table_gap))
