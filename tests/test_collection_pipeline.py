from __future__ import annotations

from hashlib import sha256
import json
from pathlib import Path
import re
import shutil
import sys
from types import ModuleType, SimpleNamespace
from typing import Mapping, Sequence

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import pytest

from pptx_wiki.collection import CollectionConfig, run_collection
from pptx_wiki.integration import (
    IntegrationConfig,
    QUALIFIED_CITATION_RE,
    validate_integrated_artifact,
)
from pptx_wiki.pipeline import PipelineConfig, run_pipeline
from pptx_wiki.quartz_publish import publish_quartz
from pptx_wiki.semantic import SemanticConfig
from pptx_wiki.source_semantic import (
    SourceIdentity,
    build_source_semantic,
    canonical_pr_number,
    extract_pr_numbers,
    load_source_semantic,
)
from pptx_wiki.wiki_output import load_provenance


_LOCAL_CITATION_RE = re.compile(r"\[slide-\d+#[^\]\s#]+\]")
_PR_RE = re.compile(r"\bPR[- ]\d+\b")


def _write_pptx(path: Path, *, pr_number: str | None, result: str) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    values = [
        f"{pr_number} reliability evaluation" if pr_number else "Reliability evaluation",
        f"Result: {result}",
        "Authoring guide: summarize the result.",
        "Authoring guide: summarize the result.",
    ]
    for index, value in enumerate(values):
        shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5 + index), Inches(7.0), Inches(0.5)
        )
        shape.name = f"TEXT_{index + 1}"
        shape.text_frame.text = value
    presentation.save(path)
    return path


def _write_prs_outside_title(
    path: Path, *, body_pr: str, table_pr: str
) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Reliability Results"

    body = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(7.0), Inches(0.8)
    )
    body.name = "REQUEST_BODY"
    body.text_frame.text = f"PR 번호: {body_pr}\nResult: passed"

    table_shape = slide.shapes.add_table(
        2, 2, Inches(0.5), Inches(2.75), Inches(7.0), Inches(1.5)
    )
    table_shape.name = "REQUEST_TABLE"
    table = table_shape.table
    table.cell(0, 0).text = "PR 번호"
    table.cell(0, 1).text = "Result"
    table.cell(1, 0).text = table_pr
    table.cell(1, 1).text = "review"

    presentation.save(path)
    return path


def _write_case_variant_identifier_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "LOT NO."

    body = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(7.0), Inches(0.8)
    )
    body.name = "REQUEST_BODY"
    body.text_frame.text = "PR-00123 reliability result\nLot No. ABC123"

    presentation.save(path)
    return path


def _write_repeated_prs_in_different_order(path: Path) -> Path:
    """Put a source-level PR before a later citation that reverses the order."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Reliability Results"

    first = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(7.0), Inches(0.8)
    )
    first.name = "FIRST_REQUEST"
    first.text_frame.text = "PR-00123 reliability result: passed"

    repeated = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(7.0), Inches(0.8)
    )
    repeated.name = "REVERSED_REQUESTS"
    repeated.text_frame.text = (
        "PR 00456 comparison result: review\n"
        "PR 00123 baseline result: passed"
    )

    presentation.save(path)
    return path


def _write_pptx_with_image(path: Path, image_path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(7.0), Inches(0.8)
    )
    text.name = "REQUEST_TEXT"
    text.text_frame.text = "PR-00123 reliability result: passed"

    Image.new("RGB", (48, 32), "navy").save(image_path, format="PNG")
    picture = slide.shapes.add_picture(
        str(image_path), Inches(0.5), Inches(1.75), width=Inches(2.0)
    )
    picture.name = "RESULT_SCREENSHOT"
    presentation.save(path)
    return path


def _matches(pattern: re.Pattern[str], value: str) -> tuple[str, ...]:
    return tuple(dict.fromkeys(match.group(0) for match in pattern.finditer(value)))


def _allowed_line(prompt: str) -> str:
    return next(
        line for line in prompt.splitlines() if line.startswith("Allowed citations:")
    )


class _ScriptedBackend:
    model = "scripted-collection-model"

    def __init__(self, *, mutate_source_pr: bool = False) -> None:
        self.mutate_source_pr = mutate_source_pr
        self.calls: list[str] = []

    def complete(
        self,
        messages: Sequence[Mapping[str, str]],
        *,
        max_tokens: int,
        temperature: float,
    ) -> str:
        prompt = messages[-1]["content"]
        if "COLLECTION_ENTITY_TOPIC_DISCOVERY" in prompt:
            self.calls.append("integration-discovery")
            citation = _matches(QUALIFIED_CITATION_RE, _allowed_line(prompt))[0]
            pr_number = _matches(_PR_RE, prompt)[0]
            return json.dumps(
                {
                    "entities": [
                        {
                            "name": pr_number,
                            "type": "project",
                            "description": f"{pr_number} source record {citation}",
                            "aliases": [],
                            "citations": [citation],
                        }
                    ],
                    "topics": [
                        {
                            "title": f"{pr_number} results",
                            "description": "Reliability result",
                            "citations": [citation],
                        }
                    ],
                }
            )
        if "COLLECTION_GROUNDED_PAGE" in prompt:
            self.calls.append("integration-page")
            citation = _matches(QUALIFIED_CITATION_RE, _allowed_line(prompt))[0]
            pr_number = _matches(_PR_RE, prompt)[0]
            return f"- {pr_number} result is recorded. {citation}"
        if "Organize this evidence" in prompt and "Return JSON only" in prompt:
            self.calls.append("source-discovery")
            evidence = prompt.split("<evidence>", 1)[1].split("</evidence>", 1)[0]
            citation = _matches(_LOCAL_CITATION_RE, evidence)[0]
            return json.dumps(
                {
                    "topics": [
                        {
                            "title": "Reliability result",
                            "description": "Relevant result only",
                            "citations": [citation],
                        }
                    ]
                }
            )
        if "Write one concise wiki page" in prompt:
            self.calls.append("source-page")
            citation = _matches(_LOCAL_CITATION_RE, _allowed_line(prompt))[0]
            pr_number = _matches(_PR_RE, prompt)[0]
            if self.mutate_source_pr:
                pr_number = f"{pr_number}-A"
            return f"# Reliability result\n\n- {pr_number} result is recorded. {citation}"
        raise AssertionError(f"unexpected backend prompt: {prompt[:120]}")


class _EntityRetryBackend(_ScriptedBackend):
    """Return a filtered PR entity first, then a grounded non-PR entity."""

    def complete(
        self,
        messages: Sequence[Mapping[str, str]],
        *,
        max_tokens: int,
        temperature: float,
    ) -> str:
        prompt = messages[-1]["content"]
        if "COLLECTION_ENTITY_ONLY_RETRY" in prompt:
            self.calls.append("integration-entity-retry")
            citation = _matches(QUALIFIED_CITATION_RE, _allowed_line(prompt))[0]
            return json.dumps(
                {
                    "entities": [
                        {
                            "name": "Reliability evaluation",
                            "type": "concept",
                            "description": (
                                f"Reliability evaluation source concept {citation}"
                            ),
                            "aliases": [],
                            "citations": [citation],
                        }
                    ]
                }
            )
        if "Write one concise wiki page" in prompt:
            self.calls.append("source-page")
            citation = _matches(_LOCAL_CITATION_RE, _allowed_line(prompt))[0]
            pr_number = _matches(_PR_RE, prompt)[0]
            return (
                "# Reliability evaluation\n\n"
                f"- {pr_number} reliability evaluation passed. {citation}"
            )
        return super().complete(
            messages,
            max_tokens=max_tokens,
            temperature=temperature,
        )


class _EmptyEntityRetryBackend(_EntityRetryBackend):
    def complete(
        self,
        messages: Sequence[Mapping[str, str]],
        *,
        max_tokens: int,
        temperature: float,
    ) -> str:
        prompt = messages[-1]["content"]
        if "COLLECTION_ENTITY_ONLY_RETRY" in prompt:
            self.calls.append("integration-entity-retry")
            return json.dumps({"entities": []})
        return super().complete(
            messages,
            max_tokens=max_tokens,
            temperature=temperature,
        )


def _install_fake_quartz(monkeypatch: pytest.MonkeyPatch) -> None:
    module = ModuleType("pptx_wiki.quartz_publish")

    def publish_quartz(
        collection_root: str | Path,
        integrated_dir: str | Path,
        output_dir: str | Path,
        *,
        site_title: str,
    ) -> SimpleNamespace:
        integrated = validate_integrated_artifact(integrated_dir)
        destination = Path(output_dir)
        content = destination / "content"
        content.mkdir(parents=True)
        (content / "index.md").write_text(
            f"---\ntitle: {site_title}\n---\n\n# {site_title}\n",
            encoding="utf-8",
            newline="\n",
        )
        manifest_path = destination / "manifest.json"
        manifest_path.write_text(
            json.dumps(
                {
                    "schema_version": "test.quartz.v1",
                    "page_count": len(integrated["pages.jsonl"]),
                },
                sort_keys=True,
            )
            + "\n",
            encoding="utf-8",
            newline="\n",
        )
        return SimpleNamespace(
            manifest_path=manifest_path,
            content_dir=content,
            page_count=len(integrated["pages.jsonl"]),
        )

    module.publish_quartz = publish_quartz  # type: ignore[attr-defined]
    monkeypatch.setitem(sys.modules, "pptx_wiki.quartz_publish", module)


@pytest.fixture()
def collection_result(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    first = _write_pptx(
        tmp_path / "first.pptx", pr_number="PR-00123", result="passed"
    )
    second = _write_pptx(
        tmp_path / "second.pptx", pr_number="PR-00456", result="review"
    )
    duplicate = tmp_path / "first-copy.pptx"
    shutil.copyfile(first, duplicate)
    backend = _ScriptedBackend()
    _install_fake_quartz(monkeypatch)
    result = run_collection(
        [first, second, duplicate],
        tmp_path / "collection",
        semantic_backend=backend,
        integration_backend=backend,
        config=CollectionConfig(
            semantic=SemanticConfig(
                goal="Keep PR-specific reliability results.",
                coverage_policy="selected",
                discover_topics=True,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(
                goal="Integrate PR-specific reliability results.",
                repair_attempts=0,
            ),
            site_title="Reliability Wiki",
        ),
    )
    return result, backend


def test_run_collection_preserves_prs_qualifies_citations_and_covers_decisions(
    collection_result,
) -> None:
    result, backend = collection_result

    assert result.input_count == 3
    assert result.unique_source_count == 2
    assert set(result.pr_numbers) == {"PR-00123", "PR-00456"}
    duplicate_source = next(
        source for source in result.sources if source.pr_numbers == ("PR-00123",)
    )
    assert len(duplicate_source.occurrences) == 2

    for source in result.sources:
        assert source.parsed.parsed_manifest_path.is_file()
        assert source.semantic.markdown_path.name == "semantic.md"
        assert source.semantic.markdown_path.is_file()
        semantic_markdown = source.semantic.markdown_path.read_text(encoding="utf-8")
        assert source.pr_numbers[0] in semantic_markdown

        loaded = load_source_semantic(source.semantic.output_dir)
        provenance = load_provenance(source.parsed.corpus.provenance_path)
        decision_citations = {item["citation"] for item in loaded["decisions"]}
        assert decision_citations == {item["citation"] for item in provenance}
        assert len(loaded["decisions"]) == len(provenance)
        assert any(
            item["disposition"] == "duplicate"
            and item["reason_code"] == "exact_duplicate"
            for item in loaded["decisions"]
        )

    integrated = validate_integrated_artifact(result.integrated.output_dir)
    source_map = integrated["source-map.jsonl"]
    by_local: dict[str, set[str]] = {}
    for item in source_map:
        by_local.setdefault(item["local_citation"], set()).add(
            item["qualified_citation"]
        )
    assert any(len(qualified) == 2 for qualified in by_local.values())
    assert len({item["qualified_citation"] for item in source_map}) == len(source_map)

    coverage = integrated["coverage.jsonl"]
    assert {item["source_id"] for item in coverage} == {
        source.source_id for source in result.sources
    }
    assert all(item["covered"] and item["page_ids"] for item in coverage)
    integrated_pages = "\n".join(
        item["body_markdown"] for item in integrated["pages.jsonl"]
    )
    assert "PR-00123" in integrated_pages
    assert "PR-00456" in integrated_pages
    assert "source-discovery" in backend.calls
    assert "integration-discovery" in backend.calls


def test_integrated_manifest_tamper_is_rejected(collection_result) -> None:
    result, _ = collection_result
    manifest_path = result.integrated.manifest_path
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["files"]["pages.jsonl"]["sha256"] = "0" * 64
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, sort_keys=True, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )

    with pytest.raises(ValueError, match=r"SHA-256 mismatch: pages\.jsonl"):
        validate_integrated_artifact(result.integrated.output_dir)


def test_source_semantic_uses_fallback_when_model_mutates_pr(
    tmp_path: Path,
) -> None:
    source = _write_pptx(
        tmp_path / "mutated.pptx", pr_number="PR-00123", result="passed"
    )
    source_root = tmp_path / "source"
    parsed = run_pipeline(source, source_root)
    backend = _ScriptedBackend(mutate_source_pr=True)

    semantic = build_source_semantic(
        parsed.corpus.output_dir,
        identity=SourceIdentity(
            source_id="deck-mutated-pr",
            source_name=source.name,
            source_sha256=sha256(source.read_bytes()).hexdigest(),
            pr_numbers=("PR-00123",),
        ),
        backend=backend,
        output_dir=source_root / "semantic",
        config=SemanticConfig(
            goal="Keep the exact PR number.",
            coverage_policy="selected",
            discover_topics=False,
            repair_attempts=0,
        ),
    )

    assert semantic.fallback is True
    markdown = semantic.markdown_path.read_text(encoding="utf-8")
    assert "PR-00123" in markdown
    assert "PR-00123-A" not in markdown
    assert any(
        "changed or invented PR number" in warning for warning in semantic.warnings
    )


def test_body_and_table_prs_reach_integrated_coverage_and_real_quartz_pages(
    tmp_path: Path,
) -> None:
    source = _write_prs_outside_title(
        tmp_path / "reliability-results.pptx",
        body_pr="PR-00123",
        table_pr="PR-00456",
    )
    assert _PR_RE.search(source.name) is None
    assert _PR_RE.search(
        Presentation(source).slides[0].shapes.title.text
    ) is None

    backend = _ScriptedBackend()
    result = run_collection(
        [source],
        tmp_path / "collection",
        semantic_backend=backend,
        integration_backend=backend,
        config=CollectionConfig(
            semantic=SemanticConfig(
                goal="Keep every exact PR-specific reliability result.",
                coverage_policy="selected",
                discover_topics=True,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(
                goal="Integrate every exact PR-specific reliability result.",
                repair_attempts=0,
            ),
            site_title="Multi-PR Reliability Wiki",
        ),
    )

    assert result.unique_source_count == 1
    assert result.pr_numbers == ("PR-00123", "PR-00456")
    assert result.sources[0].pr_numbers == ("PR-00123", "PR-00456")
    semantic_markdown = result.sources[0].semantic.markdown_path.read_text(
        encoding="utf-8"
    )
    assert "PR-00123" in semantic_markdown
    assert "PR-00456" in semantic_markdown

    integrated = validate_integrated_artifact(result.integrated.output_dir)
    assert len(integrated["coverage.jsonl"]) == 1
    coverage = integrated["coverage.jsonl"][0]
    assert coverage["source_id"] == result.sources[0].source_id
    assert coverage["semantic_document_id"] == result.sources[0].source_id
    assert coverage["pr_numbers"] == ["PR-00123", "PR-00456"]
    assert coverage["covered"] is True
    assert coverage["page_ids"]

    quartz_manifest = json.loads(
        result.quartz.manifest_path.read_text(encoding="utf-8")
    )
    assert result.quartz.pr_count == 2
    pr_pages = {
        item["pr_number"]: result.quartz.output_dir / item["page"]
        for item in quartz_manifest["prs"]
    }
    assert set(pr_pages) == {"PR-00123", "PR-00456"}
    for pr_number, page_path in pr_pages.items():
        assert page_path.is_file()
        assert pr_number in page_path.read_text(encoding="utf-8")


def test_same_canonical_pr_variants_share_inventory_but_remain_in_audit(
    tmp_path: Path,
) -> None:
    source = _write_prs_outside_title(
        tmp_path / "variant-spellings.pptx",
        body_pr="PR-00123",
        table_pr="PR 00123",
    )
    source_root = tmp_path / "source"
    parsed = run_pipeline(source, source_root)

    semantic = build_source_semantic(
        parsed.corpus.output_dir,
        identity=SourceIdentity(
            source_id="deck-pr-variant",
            source_name=source.name,
            source_sha256=sha256(source.read_bytes()).hexdigest(),
            pr_numbers=("PR-00123",),
        ),
        backend=_ScriptedBackend(),
        output_dir=source_root / "semantic",
        config=SemanticConfig(
            goal="Keep exact PR spellings in the audit trail.",
            coverage_policy="selected",
            discover_topics=False,
            repair_attempts=0,
        ),
    )

    loaded = load_source_semantic(semantic.output_dir)
    assert semantic.pr_numbers == ("PR-00123",)
    assert loaded["manifest"]["source_identity"]["pr_numbers"] == ["PR-00123"]
    ledger = loaded["manifest"]["pr_ledger"]
    assert {item["value"] for item in ledger} == {"PR-00123", "PR 00123"}
    assert {item["canonical_key"] for item in ledger} == {
        canonical_pr_number("PR-00123")
    }
    assert canonical_pr_number("PR-00123") == canonical_pr_number("PR 00123")


def test_korean_request_number_table_label_is_a_pr_identifier() -> None:
    markdown = """| 의뢰번호 | 결과 |
| --- | --- |
| EF236_123_4321 | PASS |
| EF236_123_9876 | PASS |
"""

    assert extract_pr_numbers(markdown) == (
        "EF236_123_4321",
        "EF236_123_9876",
    )

    merged_html = """<table>
<tr><td>의뢰번호</td><td>항목</td><td>Lot ID</td></tr>
<tr><td rowspan="3">EF236_123_4321</td><td>Precond</td><td>LEG1_AB203DLE</td></tr>
<tr><td rowspan="2">박리가혹</td><td>LEG1_AB203DLE</td></tr>
<tr><td>LEG2_AB203DLQ</td></tr>
</table>"""
    assert extract_pr_numbers(merged_html) == ("EF236_123_4321",)


def test_integrated_identifier_tokens_do_not_cross_line_boundaries(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = _write_prs_outside_title(
        tmp_path / "identifier-boundary.pptx",
        body_pr="PR-00123",
        table_pr="PR-00456",
    )
    _install_fake_quartz(monkeypatch)
    result = run_collection(
        [source],
        tmp_path / "collection",
        semantic_backend=_ScriptedBackend(),
        integration_backend=_ScriptedBackend(),
        config=CollectionConfig(
            semantic=SemanticConfig(
                coverage_policy="selected",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(repair_attempts=0),
        ),
    )

    integrated = validate_integrated_artifact(result.integrated.output_dir)
    tokens = [
        token
        for record in integrated["source-map.jsonl"]
        for token in record["identifier_tokens"]
    ]
    assert tokens
    assert all("\r" not in token and "\n" not in token for token in tokens)


def test_integration_retries_llm_when_initial_entity_candidates_filter_out(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = _write_pptx(
        tmp_path / "fallback-entity.pptx",
        pr_number="PR-00123",
        result="passed",
    )
    _install_fake_quartz(monkeypatch)
    backend = _EntityRetryBackend()
    result = run_collection(
        [source],
        tmp_path / "entity-retry",
        semantic_backend=backend,
        integration_backend=backend,
        config=CollectionConfig(
            semantic=SemanticConfig(
                coverage_policy="selected",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(repair_attempts=0),
        ),
    )
    entities = validate_integrated_artifact(result.integrated.output_dir)[
        "entities.jsonl"
    ]

    # Initial discovery proposes only the protected PR identifier. The second
    # LLM call must be entity-only and accept an exact-visible non-PR concept.
    assert backend.calls.count("integration-discovery") == 1
    assert backend.calls.count("integration-entity-retry") == 1
    entity = next(
        entity
        for entity in entities
        if entity["canonical_name"].casefold() == "reliability evaluation"
    )
    assert entity["type"] == "concept"
    assert entity["citations"]
    assert entity["source_ids"] == [result.sources[0].source_id]
    assert entity["pr_numbers"] == ["PR-00123"]

    published = publish_quartz(
        result.output_dir,
        result.integrated.output_dir,
        tmp_path / "entity-retry-quartz",
        site_title="Entity Retry Test Wiki",
    )
    assert published.entity_count == 1
    entity_page = published.content_dir / "entities" / f"{entity['id']}.md"
    assert entity_page.is_file()
    assert "Reliability evaluation" in entity_page.read_text(encoding="utf-8")


def test_integration_warns_when_entity_only_llm_retry_is_still_empty(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = _write_pptx(
        tmp_path / "empty-entity-retry.pptx",
        pr_number="PR-00123",
        result="passed",
    )
    _install_fake_quartz(monkeypatch)
    backend = _EmptyEntityRetryBackend()

    result = run_collection(
        [source],
        tmp_path / "empty-entity-retry",
        semantic_backend=backend,
        integration_backend=backend,
        config=CollectionConfig(
            semantic=SemanticConfig(
                coverage_policy="selected",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(repair_attempts=0),
        ),
    )

    assert backend.calls.count("integration-entity-retry") == 1
    assert result.integrated.entity_count == 0
    assert any(
        "entity" in warning.casefold()
        and any(word in warning.casefold() for word in ("empty", "retry", "usable"))
        for warning in result.integrated.warnings
    )


def test_integrated_identifier_tokens_deduplicate_case_variants(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = _write_case_variant_identifier_pptx(
        tmp_path / "identifier-case-variant.pptx"
    )
    _install_fake_quartz(monkeypatch)

    result = run_collection(
        [source],
        tmp_path / "collection",
        semantic_backend=_ScriptedBackend(),
        integration_backend=_ScriptedBackend(),
        config=CollectionConfig(
            semantic=SemanticConfig(
                coverage_policy="selected",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(repair_attempts=0),
        ),
    )

    source_map_path = result.integrated.source_map_path
    source_map = [
        json.loads(line)
        for line in source_map_path.read_text(encoding="utf-8").splitlines()
        if line.strip()
    ]
    legacy_record = next(
        record
        for record in source_map
        if "LOT NO." in record["identifier_tokens"]
    )
    legacy_record["identifier_tokens"].append("Lot No.")
    source_map_text = "".join(
        json.dumps(record, ensure_ascii=False, sort_keys=True) + "\n"
        for record in source_map
    )
    source_map_path.write_text(
        source_map_text,
        encoding="utf-8",
        newline="\n",
    )
    manifest_path = result.integrated.manifest_path
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["files"]["source-map.jsonl"]["sha256"] = sha256(
        source_map_text.encode("utf-8")
    ).hexdigest()
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, sort_keys=True, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )

    published = publish_quartz(
        result.output_dir,
        result.integrated.output_dir,
        tmp_path / "republished-quartz",
        site_title="Identifier Deduplication Test Wiki",
    )
    assert published.manifest_path.is_file()
    assert legacy_record["identifier_tokens"][-2:] == ["LOT NO.", "Lot No."]


def test_quartz_resume_accepts_citation_pr_order_emitted_by_integration(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = _write_repeated_prs_in_different_order(
        tmp_path / "repeated-pr-order.pptx"
    )
    _install_fake_quartz(monkeypatch)

    result = run_collection(
        [source],
        tmp_path / "collection",
        semantic_backend=_ScriptedBackend(),
        integration_backend=_ScriptedBackend(),
        config=CollectionConfig(
            semantic=SemanticConfig(
                coverage_policy="selected",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(repair_attempts=0),
        ),
    )

    source_map = [
        json.loads(line)
        for line in result.integrated.source_map_path.read_text(
            encoding="utf-8"
        ).splitlines()
        if line.strip()
    ]
    reversed_record = next(
        record
        for record in source_map
        if record["pr_variants"] == ["PR 00456", "PR 00123"]
    )
    assert result.sources[0].pr_numbers == ("PR-00123", "PR 00456")
    assert reversed_record["pr_numbers"] == ["PR 00456", "PR-00123"]

    published = publish_quartz(
        result.output_dir,
        result.integrated.output_dir,
        tmp_path / "resumed-quartz",
        site_title="Repeated PR Order Test Wiki",
    )
    assert published.manifest_path.is_file()


def test_default_collection_omits_pictures_from_parsed_semantic_and_quartz(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = _write_pptx_with_image(
        tmp_path / "ignored-picture.pptx", tmp_path / "source-image.png"
    )
    _install_fake_quartz(monkeypatch)
    backend = _ScriptedBackend()
    result = run_collection(
        [source],
        tmp_path / "collection-without-pictures",
        semantic_backend=backend,
        integration_backend=backend,
        config=CollectionConfig(
            semantic=SemanticConfig(
                coverage_policy="complete",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(repair_attempts=0),
        ),
    )

    parsed = result.sources[0].parsed
    provenance = load_provenance(parsed.corpus.provenance_path)
    assert not any(
        record["kind"] in {"image", "picture", "ocr_table", "ocr_text"}
        or record["asset_path"]
        for record in provenance
    )
    assert not list((parsed.parsed_dir / "source-assets" / "images").glob("*"))
    parsed_markdown = "\n".join(
        path.read_text(encoding="utf-8") for path in parsed.corpus.slide_paths
    )
    semantic_markdown = result.sources[0].semantic.markdown_path.read_text(
        encoding="utf-8"
    )
    assert "![" not in parsed_markdown
    assert "<img" not in parsed_markdown.casefold()
    assert "![" not in semantic_markdown
    assert "<img" not in semantic_markdown.casefold()

    published = publish_quartz(
        result.output_dir,
        result.integrated.output_dir,
        tmp_path / "quartz-without-pictures",
        site_title="Text-only Test Wiki",
    )
    assert published.asset_paths == ()
    assert not (published.content_dir / "assets").exists()
    quartz_manifest = json.loads(
        published.manifest_path.read_text(encoding="utf-8")
    )
    assert quartz_manifest["asset_count"] == 0
    for markdown_path in published.content_dir.rglob("*.md"):
        markdown = markdown_path.read_text(encoding="utf-8")
        assert "![" not in markdown
        assert "<img" not in markdown.casefold()


def test_quartz_rejects_tampered_collection_image_asset(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = _write_pptx_with_image(
        tmp_path / "asset-lineage.pptx", tmp_path / "source-image.png"
    )
    _install_fake_quartz(monkeypatch)
    backend = _ScriptedBackend()
    result = run_collection(
        [source],
        tmp_path / "collection",
        semantic_backend=backend,
        integration_backend=backend,
        config=CollectionConfig(
            pipeline=PipelineConfig(include_images=True),
            semantic=SemanticConfig(
                coverage_policy="complete",
                discover_topics=False,
                repair_attempts=0,
            ),
            integration=IntegrationConfig(repair_attempts=0),
        ),
    )

    image_assets = list(
        (result.sources[0].parsed.parsed_dir / "source-assets" / "images").glob("*")
    )
    assert len(image_assets) == 1
    asset = image_assets[0]
    provenance = load_provenance(result.sources[0].parsed.corpus.provenance_path)
    image_record = next(record for record in provenance if record["asset_path"])
    assert sha256(asset.read_bytes()).hexdigest() == image_record["metadata"][
        "image_sha256"
    ]

    asset.write_bytes(asset.read_bytes() + b"tampered")
    with pytest.raises(ValueError, match=r"(?i)asset.*(?:SHA-256|digest|changed)"):
        publish_quartz(
            result.output_dir,
            result.integrated.output_dir,
            tmp_path / "republished-quartz",
            site_title="Tamper Test Wiki",
        )


def test_collection_without_pr_fails_closed_before_llm(tmp_path: Path) -> None:
    source = _write_pptx(
        tmp_path / "missing-pr.pptx", pr_number=None, result="passed"
    )
    backend = _ScriptedBackend()
    output = tmp_path / "collection"

    with pytest.raises(ValueError, match="PR"):
        run_collection(
            [source],
            output,
            semantic_backend=backend,
            integration_backend=backend,
            config=CollectionConfig(
                semantic=SemanticConfig(
                    goal="Keep PR-specific results.",
                    coverage_policy="selected",
                    repair_attempts=0,
                ),
                integration=IntegrationConfig(repair_attempts=0),
            ),
        )

    assert backend.calls == []
    manifest = json.loads(
        (output / "collection-manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["status"] == "failed"
    assert manifest["error"]["type"] == "ValueError"
    assert not (output / ".collection.lock").exists()
