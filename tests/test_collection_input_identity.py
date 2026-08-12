from __future__ import annotations

import os
from hashlib import sha256
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
import pytest

from pptx_wiki.collection import (
    _stage_verified_snapshot,
    _verify_current_path_digest,
    discover_pptx_inputs,
)


def _write_minimal_pptx(path: Path, label: str) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(7.0), Inches(0.8)
    )
    textbox.text_frame.text = f"PR-00123 {label}"
    presentation.save(path)
    return path


def _transient_samestat_mismatch(monkeypatch: pytest.MonkeyPatch) -> list[bool]:
    real_samestat = os.path.samestat
    results: list[bool] = []

    def samestat(left: os.stat_result, right: os.stat_result) -> bool:
        if not results:
            results.append(False)
            return False
        result = real_samestat(left, right)
        results.append(result)
        return result

    monkeypatch.setattr(os.path, "samestat", samestat)
    return results


def test_directory_discovery_retries_a_transient_open_identity_mismatch(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """A transient Windows/virtual-FS file-ID read must not break folder runs."""

    input_dir = tmp_path / "input"
    input_dir.mkdir()
    first = _write_minimal_pptx(input_dir / "first.pptx", "first")
    second = _write_minimal_pptx(input_dir / "second.pptx", "second")
    results = _transient_samestat_mismatch(monkeypatch)

    occurrences = discover_pptx_inputs([input_dir])

    assert {item.path for item in occurrences} == {first.resolve(), second.resolve()}
    assert results[:2] == [False, True]


def test_snapshot_staging_retries_a_transient_open_identity_mismatch(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """The second open used by a collection run needs the same retry behavior."""

    source = _write_minimal_pptx(tmp_path / "source.pptx", "staging")
    occurrence = discover_pptx_inputs([source])[0]
    results = _transient_samestat_mismatch(monkeypatch)

    snapshot = _stage_verified_snapshot(
        occurrence,
        tmp_path / "snapshot",
        reject_external_resources=True,
    )

    assert snapshot.read_bytes() == source.read_bytes()
    assert results[:2] == [False, True]


def test_directory_discovery_uses_hash_fallback_for_persistent_identity_mismatch(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """Providers with permanently unreliable file IDs can use stable bytes."""

    input_dir = tmp_path / "input"
    input_dir.mkdir()
    source = _write_minimal_pptx(input_dir / "source.pptx", "persistent mismatch")
    monkeypatch.setattr(os.path, "samestat", lambda _left, _right: False)

    occurrences = discover_pptx_inputs([input_dir])

    assert [item.path for item in occurrences] == [source.resolve()]
    assert occurrences[0].sha256 == sha256(source.read_bytes()).hexdigest()


def test_identity_fallback_rejects_different_current_bytes(tmp_path: Path) -> None:
    source = _write_minimal_pptx(tmp_path / "source.pptx", "before")
    expected = sha256(source.read_bytes()).hexdigest()
    _write_minimal_pptx(source, "after replacement")

    with pytest.raises(ValueError, match="input path changed while opening"):
        _verify_current_path_digest(
            source,
            expected,
            error_message="input path changed while opening: source.pptx",
        )
